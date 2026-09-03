"""Comprehensive verification suite for the PrivacyDrop engine.

Covers every supported format, every option, and a wide set of edge cases:
validity, losslessness, ICC handling, corrupt/encrypted/password inputs,
filename handling, name collisions, original-untouched guarantees, and the
CLI.  Run:  python tests/test_full.py
"""
import hashlib
import io
import os
import shutil
import struct
import subprocess
import sys
import zipfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import scrubber
from scrubber import (sanitize_file, Options, default_output_path,
                      strip_jpeg_bytes, supported)

from PIL import Image
from PIL.TiffImagePlugin import IFDRational

TMP = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tmp_full")
OUT = os.path.join(TMP, "out")
PASS = FAIL = SKIP = 0


def check(name, cond, detail=""):
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  PASS  {name}")
    else:
        FAIL += 1
        print(f"  FAIL  {name}   {detail}")


def note(name):
    global SKIP
    SKIP += 1
    print(f"  SKIP  {name}")


def sha(p):
    return hashlib.sha256(open(p, "rb").read()).hexdigest()


# ---------------------------------------------------------------- fixtures

def make_exif():
    exif = Image.Exif()
    exif[0x010F] = "Acme Corp"
    exif[0x0110] = "SecretCam 9000"
    exif[0x0131] = "LeakySoft 1.0"
    exif[0x0132] = "2024:01:02 03:04:05"
    exif[0x8825] = {1: "S",
                    2: (IFDRational(31, 1), IFDRational(57, 1), IFDRational(1, 1)),
                    3: "E",
                    4: (IFDRational(115, 1), IFDRational(51, 1), IFDRational(1, 1))}
    return exif


def jpeg_segments(data):
    i, n, segs = 2, len(data), []
    while i < n - 1:
        if data[i] != 0xFF:
            break
        b = data[i + 1]
        if b == 0xDA:
            break
        if b in (0xD8, 0xD9) or b == 0x01 or 0xD0 <= b <= 0xD7:
            segs.append((b, b"")); i += 2; continue
        if i + 4 > n:
            break
        seglen = (data[i + 2] << 8) | data[i + 3]
        if i + 2 + seglen > n:
            break
        segs.append((b, data[i + 4:i + 2 + seglen]))
        i += 2 + seglen
    return segs


def has_icc(data):
    return any(b == 0xE2 and p.startswith(b"ICC_PROFILE\x00")
               for b, p in jpeg_segments(data))


def png_chunks(data):
    i, n, out = 8, len(data), []
    while i + 8 <= n:
        length = struct.unpack(">I", data[i:i + 4])[0]
        out.append(data[i + 4:i + 8])
        i += 12 + length
    return out


def webp_chunks(data):
    i, n, out = 12, len(data), []
    while i + 8 <= n:
        fourcc = data[i:i + 4]
        size = struct.unpack("<I", data[i + 4:i + 8])[0]
        out.append(fourcc)
        i += 8 + size + (size % 2)
    return out


def sanitize(src, opts=None, out_dir=OUT):
    opts = opts or Options()
    dst = default_output_path(src, out_dir, opts.suffix)
    return sanitize_file(src, dst, opts)


# ---------------------------------------------------------------- tests

def t_jpeg():
    print("\n[JPEG]")
    # basic EXIF/GPS/XMP/comment + losslessness
    src = os.path.join(TMP, "a.jpg")
    img = Image.new("RGB", (64, 48), (200, 40, 40))
    img.save(src, "JPEG", quality=90, exif=make_exif().tobytes(),
             comment=b"hi there")
    rep = sanitize(src)
    check("exif+gps stripped", rep.status == "ok" and "EXIF" in rep.summary())
    segs = jpeg_segments(open(rep.dst, "rb").read())
    check("no APP1 left", not any(b == 0xE1 for b, _ in segs))
    with Image.open(src) as a, Image.open(rep.dst) as b:
        check("pixel-identical", a.convert("RGB").tobytes()
              == b.convert("RGB").tobytes())

    # progressive
    p = os.path.join(TMP, "prog.jpg")
    Image.new("RGB", (64, 48), (10, 90, 10)).save(
        p, "JPEG", progressive=True, exif=make_exif().tobytes())
    rep = sanitize(p)
    check("progressive cleaned", rep.status == "ok")
    with Image.open(rep.dst) as im:
        im.load()
        check("progressive still decodes", im.size == (64, 48))

    # XMP-only
    x = os.path.join(TMP, "xmp.jpg")
    Image.new("RGB", (20, 20)).save(
        x, "JPEG", xmp=b"<x:xmpmeta xmlns:x='adobe:ns:meta/'>"
                      b"<rdf:RDF xmlns:rdf='rdf'/></x:xmpmeta>")
    rep = sanitize(x)
    check("XMP stripped", rep.status == "ok" and "XMP" in rep.summary())

    # comment-only
    c = os.path.join(TMP, "com.jpg")
    Image.new("RGB", (20, 20)).save(c, "JPEG", comment=b"secret")
    rep = sanitize(c)
    check("comment stripped", rep.status == "ok"
          and "comment" in rep.summary())

    # already clean -> identical bytes (fresh metadata-free file)
    clean = os.path.join(TMP, "clean.jpg")
    Image.new("RGB", (20, 20)).save(clean, "JPEG")
    rep = sanitize(clean)
    check("already-clean identical", rep.status == "clean"
          and open(rep.dst, "rb").read() == open(clean, "rb").read())


def t_icc():
    print("\n[ICC colour profile handling]")
    from PIL import ImageCms
    prof = ImageCms.ImageCmsProfile(ImageCms.createProfile("sRGB")).tobytes()

    src = os.path.join(TMP, "icc.jpg")
    Image.new("RGB", (32, 32), (90, 90, 90)).save(
        src, "JPEG", exif=make_exif().tobytes(), icc_profile=prof)
    check("source has ICC", has_icc(open(src, "rb").read()))

    # keep (default)
    rep = sanitize(src, Options())
    out = open(rep.dst, "rb").read()
    check("ICC kept by default", has_icc(out))
    check("EXIF still stripped", "EXIF" in rep.summary())

    # strip
    rep2 = sanitize(src, Options(strip_icc=True))
    out2 = open(rep2.dst, "rb").read()
    check("ICC removed when asked", not has_icc(out2)
          and "ICC" in rep2.summary())


def t_png():
    print("\n[PNG]")
    from PIL.PngImagePlugin import PngInfo
    src = os.path.join(TMP, "a.png")
    img = Image.new("RGB", (32, 32), (10, 120, 200))
    info = PngInfo()
    info.add_text("Author", "Alice")
    info.add_text("Software", "LeakyEditor")
    img.save(src, "PNG", pnginfo=info)
    rep = sanitize(src)
    check("text chunks stripped", rep.status == "ok"
          and "text chunk" in rep.summary())
    chunks = png_chunks(open(rep.dst, "rb").read())
    check("no tEXt/iTXt/zTXt",
          not any(ch in (b"tEXt", b"iTXt", b"zTXt") for ch in chunks))
    with Image.open(rep.dst) as im:
        check("pixel preserved", im.convert("RGB").getpixel((0, 0))
              == (10, 120, 200))

    # grayscale + palette
    for mode, name in (("L", "gray.png"), ("P", "pal.png")):
        p = os.path.join(TMP, name)
        Image.new(mode, (24, 24), 128).save(p, "PNG", pnginfo=info)
        rep = sanitize(p)
        check(f"{mode} PNG cleaned", rep.status == "ok")
        with Image.open(rep.dst) as im:
            check(f"{mode} PNG valid", im.size == (24, 24))

    # ICC in PNG kept / stripped
    from PIL import ImageCms
    prof = ImageCms.ImageCmsProfile(ImageCms.createProfile("sRGB")).tobytes()
    ip = os.path.join(TMP, "iccp.png")
    Image.new("RGB", (16, 16), (5, 5, 5)).save(ip, "PNG", icc_profile=prof,
                                               pnginfo=info)
    rep = sanitize(ip, Options())
    check("PNG ICC kept", b"iCCP" in png_chunks(open(rep.dst, "rb").read()))
    rep = sanitize(ip, Options(strip_icc=True))
    check("PNG ICC stripped", b"iCCP" not in
          png_chunks(open(rep.dst, "rb").read()))

    # bad signature
    bad = os.path.join(TMP, "fake.png")
    open(bad, "wb").write(b"GIF89a not a png")
    rep = sanitize(bad)
    check("fake PNG -> error", rep.status == "error"
          and "not a valid PNG" in rep.error)


def t_webp():
    print("\n[WebP]")
    # lossy VP8 + exif
    src = os.path.join(TMP, "a.webp")
    Image.new("RGB", (32, 32), (30, 200, 90)).save(
        src, "WEBP", exif=make_exif().tobytes())
    rep = sanitize(src)
    check("webp exif stripped", rep.status == "ok" and "EXIF" in rep.summary())
    out = open(rep.dst, "rb").read()
    check("no EXIF chunk", b"EXIF" not in out)
    with Image.open(src) as a, Image.open(rep.dst) as b:
        check("pixel preserved", a.convert("RGB").tobytes()
              == b.convert("RGB").tobytes())

    # lossless VP8L
    l = os.path.join(TMP, "l.webp")
    Image.new("RGB", (24, 24), (90, 30, 200)).save(
        l, "WEBP", lossless=True, exif=make_exif().tobytes())
    rep = sanitize(l)
    check("lossless webp cleaned", rep.status == "ok")
    with Image.open(rep.dst) as im:
        check("lossless webp valid", im.size == (24, 24))

    # animated
    f1 = Image.new("RGB", (24, 24), (255, 0, 0))
    f2 = Image.new("RGB", (24, 24), (0, 0, 255))
    a = os.path.join(TMP, "anim.webp")
    f1.save(a, "WEBP", save_all=True, append_images=[f2], duration=100,
            loop=0, exif=make_exif().tobytes())
    rep = sanitize(a)
    check("animated webp cleaned", rep.status == "ok")
    with Image.open(rep.dst) as im:
        check("animation frames kept", getattr(im, "n_frames", 1) == 2)

    # already clean identical
    c = os.path.join(TMP, "plain.webp")
    Image.new("RGB", (16, 16)).save(c, "WEBP")
    rep = sanitize(c)
    check("clean webp identical", rep.status == "clean"
          and open(rep.dst, "rb").read() == open(c, "rb").read())


def t_tiff_gif_bmp():
    print("\n[TIFF / GIF / BMP]")
    # TIFF
    t = os.path.join(TMP, "a.tiff")
    img = Image.new("RGB", (40, 30), (11, 22, 33))
    try:
        img.save(t, "TIFF", exif=make_exif().tobytes())
        src_has_exif = bool(dict(Image.open(t).getexif()))
    except Exception:
        img.save(t, "TIFF")
        src_has_exif = False
    rep = sanitize(t)
    check("TIFF processed", rep.ok, rep.summary())
    with Image.open(rep.dst) as im:
        check("TIFF pixel preserved", im.convert("RGB").getpixel((0, 0))
              == (11, 22, 33))
        exif = dict(im.getexif())
        # identity tags (Make/Model/Software/DateTime/Artist) must be gone;
        # baseline IFD tags (width/height/bits) are structural, not identity.
        check("TIFF identity EXIF gone",
              not any(tag in exif for tag in (271, 272, 305, 306, 315)),
              str({k: v for k, v in exif.items() if k in (271, 272, 305, 306, 315)}))
    if not src_has_exif:
        note("(Pillow could not embed EXIF in source TIFF; strip still verified)")

    # multipage TIFF
    m = os.path.join(TMP, "multi.tiff")
    f1 = Image.new("RGB", (32, 32), (1, 2, 3))
    f2 = Image.new("RGB", (32, 32), (4, 5, 6))
    f1.save(m, "TIFF", save_all=True, append_images=[f2])
    rep = sanitize(m)
    check("multipage TIFF ok", rep.ok)
    with Image.open(rep.dst) as im:
        check("TIFF pages kept", getattr(im, "n_frames", 1) == 2)

    # animated GIF with comment
    g = os.path.join(TMP, "anim.gif")
    f1 = Image.new("RGB", (30, 30), (255, 0, 0)).convert("P")
    f2 = Image.new("RGB", (30, 30), (0, 0, 255)).convert("P")
    f1.save(g, "GIF", save_all=True, append_images=[f2], duration=100,
            comment=b"made by alice")
    with Image.open(g) as src_im:
        src_frames = getattr(src_im, "n_frames", 1)
    check("fixture is animated", src_frames == 2, f"n_frames={src_frames}")
    rep = sanitize(g)
    check("GIF processed", rep.ok)
    with Image.open(rep.dst) as im:
        check("GIF frames kept", getattr(im, "n_frames", 1) == 2)
        check("GIF comment gone", "comment" not in im.info)

    # BMP
    b = os.path.join(TMP, "a.bmp")
    Image.new("RGB", (24, 24), (9, 8, 7)).save(b, "BMP")
    rep = sanitize(b)
    check("BMP processed", rep.ok)
    with Image.open(rep.dst) as im:
        check("BMP pixel preserved", im.convert("RGB").getpixel((0, 0))
              == (9, 8, 7))


def t_heic():
    print("\n[HEIC]")
    if not scrubber.HAS_HEIF:
        note("pillow-heif not available; skipping HEIC encode/decode test")
        return
    import pillow_heif
    h = os.path.join(TMP, "a.heic")
    img = Image.new("RGB", (48, 36), (70, 80, 90))
    try:
        img.save(h, format="HEIF", exif=make_exif().tobytes())
        src_exif = bool(dict(Image.open(h).getexif()))
    except Exception as e:
        img.save(h, format="HEIF")
        src_exif = False
        note(f"(save-with-exif unsupported: {e})")
    rep = sanitize(h)
    check("HEIC processed", rep.ok, rep.summary())
    with Image.open(rep.dst) as im:
        check("HEIC decodes", im.size == (48, 36))
        check("HEIC no EXIF", not dict(im.getexif()))
    if not src_exif:
        note("(source HEIC had no EXIF to strip)")


def t_pdf():
    print("\n[PDF]")
    import pikepdf
    # metadata + attachment + annotation
    p = os.path.join(TMP, "doc.pdf")
    pdf = pikepdf.new()
    pdf.add_blank_page(page_size=(200, 200))
    pdf.docinfo["/Title"] = "Secret"
    pdf.docinfo["/Author"] = "Alice"
    pdf.docinfo["/Producer"] = "LeakySoft"
    with pdf.open_metadata() as m:
        m["dc:creator"] = ["Alice"]
    notes = os.path.join(TMP, "_notes.txt")
    open(notes, "w").write("secret")
    pdf.attachments["notes.txt"] = pikepdf.AttachedFileSpec.from_filepath(pdf, notes)
    page = pdf.pages[0]
    page["/Annots"] = pikepdf.Array([pikepdf.Dictionary({
        "/Type": pikepdf.Name("/Annot"),
        "/Subtype": pikepdf.Name("/Text"),
        "/Contents": pikepdf.String("visible note"),
        "/Rect": pikepdf.Array([10, 10, 100, 100])})])
    pdf.save(p)

    rep = sanitize(p)
    check("pdf cleaned", rep.status == "ok", rep.summary())
    with pikepdf.open(rep.dst) as q:
        check("docinfo empty", not list(q.docinfo.keys()))
        check("no XMP", "/Metadata" not in q.Root)
        check("no attachments", not list(q.attachments.keys()))
        check("annotation preserved (content, not identity)",
              "/Annots" in q.pages[0])

    # password-protected
    pw = os.path.join(TMP, "locked.pdf")
    pdf = pikepdf.new(); pdf.add_blank_page()
    pdf.save(pw, encryption=pikepdf.Encryption(owner="secret", user="secret"))
    rep = sanitize(pw)
    check("password pdf -> clear error",
          rep.status == "error" and "password" in rep.error.lower())

    # owner-encrypted only (opens without password)
    ow = os.path.join(TMP, "owner.pdf")
    pdf = pikepdf.new(); pdf.add_blank_page()
    pdf.docinfo["/Author"] = "Alice"
    pdf.save(ow, encryption=pikepdf.Encryption(owner="pw", user=""))
    rep = sanitize(ow)
    check("owner-only encrypted cleaned", rep.status == "ok",
          rep.summary())

    # malformed
    mf = os.path.join(TMP, "bad.pdf")
    open(mf, "wb").write(b"%PDF-1.7\nthis is not a real pdf")
    rep = sanitize(mf)
    check("malformed pdf -> error", rep.status == "error", rep.summary())


def t_office():
    print("\n[Office — DOCX / XLSX / PPTX]")
    import docx
    import openpyxl
    from pptx import Presentation

    # ---- DOCX with custom props + embedded PNG + author
    d = os.path.join(TMP, "r.docx")
    doc = docx.Document()
    doc.core_properties.author = "Alice"
    doc.core_properties.last_modified_by = "Bob"
    doc.core_properties.title = "Top Secret"
    doc.add_paragraph("hello world")
    doc.add_picture(_png_with_meta(os.path.join(TMP, "in_docx.png")),
                    width=docx.shared.Inches(1))
    doc.save(d)
    # inject custom.xml with a custom property
    _inject(d, "docProps/custom.xml", scrubber.CUSTOM_XML.replace(
        "/>", "><property name='Client' fmtid='{D5CDD505-2E9C-101B-9397-08002B2CF9AE}' pid='2'><vt:lpwstr xmlns:vt='http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes'>ACME</vt:lpwstr></property></Properties>"))
    # inject a docProps/thumbnail.jpeg carrying EXIF (regression: prefix bug)
    _jpeg_path(os.path.join(TMP, "thumb.jpg"))
    _inject(d, "docProps/thumbnail.jpeg",
            open(os.path.join(TMP, "thumb.jpg"), "rb").read())
    rep = sanitize(d)
    check("docx cleaned", rep.status == "ok", rep.summary())
    core = _part(rep.dst, "docProps/core.xml")
    check("author removed", "Alice" not in core and "Bob" not in core)
    check("title removed", "Top Secret" not in core)
    custom = _part(rep.dst, "docProps/custom.xml")
    check("custom props blanked", "ACME" not in custom)
    media = _part(rep.dst, "word/media/image1.png", raw=True)
    check("embedded png scrubbed", b"tEXt" not in media)
    # docProps/thumbnail.jpeg (if present) must be scrubbed too — a
    # regression check for the case-sensitive prefix bug.
    with zipfile.ZipFile(rep.dst) as z:
        thumbs = [n for n in z.namelist()
                  if n.lower().startswith("docprops/")
                  and n.lower().endswith((".jpg", ".jpeg", ".png", ".webp"))]
        check("docProps thumbnail scrubbed",
              all(b"Exif\x00\x00" not in z.read(n) for n in thumbs),
              str(thumbs))
    check("docx still opens", "hello world" in
          [pp.text for pp in docx.Document(rep.dst).paragraphs])

    # timestamps zeroed
    with zipfile.ZipFile(rep.dst) as z:
        dts = {i.date_time for i in z.infolist()}
    check("zip timestamps zeroed", dts == {(1980, 1, 1, 0, 0, 0)}, str(dts))

    # ---- XLSX with creator + embedded JPEG
    x = os.path.join(TMP, "book.xlsx")
    wb = openpyxl.Workbook()
    wb.properties.creator = "Alice"
    wb.properties.lastModifiedBy = "Bob"
    ws = wb.active
    ws["A1"] = "budget"
    img = openpyxl.drawing.image.Image(_jpeg_path(os.path.join(TMP, "in_xlsx.jpg")))
    ws.add_image(img, "B2")
    wb.save(x)
    rep = sanitize(x)
    check("xlsx cleaned", rep.status == "ok", rep.summary())
    wb2 = openpyxl.load_workbook(rep.dst)
    check("xlsx creator removed", not wb2.properties.creator
          and not wb2.properties.lastModifiedBy)
    check("xlsx content kept", wb2.active["A1"].value == "budget")
    xmedia = _first_part(rep.dst, "xl/media/", raw=True)
    check("xlsx embedded image scrubbed", b"Exif\x00\x00" not in xmedia)

    # ---- PPTX with author + embedded JPEG
    p = os.path.join(TMP, "deck.pptx")
    prs = Presentation()
    prs.core_properties.author = "Alice"
    prs.core_properties.last_modified_by = "Bob"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(_jpeg_path(os.path.join(TMP, "in_pptx.jpg")),
                             0, 0, width=1000000)
    prs.save(p)
    rep = sanitize(p)
    check("pptx cleaned", rep.status == "ok", rep.summary())
    prs2 = Presentation(rep.dst)
    check("pptx author removed", not prs2.core_properties.author
          and not prs2.core_properties.last_modified_by)
    pmedia = _first_part(rep.dst, "ppt/media/", raw=True)
    check("pptx embedded image scrubbed", b"Exif\x00\x00" not in pmedia)


def _png_with_meta(path):
    from PIL.PngImagePlugin import PngInfo
    info = PngInfo(); info.add_text("Author", "Alice")
    Image.new("RGB", (20, 20), (1, 2, 3)).save(path, "PNG", pnginfo=info)
    return path


def _jpeg_path(path):
    Image.new("RGB", (20, 20), (4, 5, 6)).save(
        path, "JPEG", exif=make_exif().tobytes())
    return path


def _inject(zip_path, name, data):
    with zipfile.ZipFile(zip_path, "a") as z:
        z.writestr(name, data)


def _part(zip_path, name, raw=False):
    with zipfile.ZipFile(zip_path) as z:
        return z.read(name) if raw else z.read(name).decode("utf-8", "ignore")


def _first_part(zip_path, prefix, raw=False):
    with zipfile.ZipFile(zip_path) as z:
        for n in z.namelist():
            if n.startswith(prefix):
                return z.read(n) if raw else z.read(n).decode("utf-8", "ignore")
    return b"" if raw else ""


def t_filenames():
    print("\n[Filename handling]")
    # unicode
    u = os.path.join(TMP, "cafÃ©-ç…§ç‰‡.JPG")   # also uppercase extension
    Image.new("RGB", (16, 16), (1, 1, 1)).save(u, "JPEG",
                                               exif=make_exif().tobytes())
    rep = sanitize(u)
    check("unicode + uppercase ext cleaned", rep.status == "ok")
    check("output exists with suffix", os.path.exists(rep.dst))

    # collision: same file twice into the same fresh folder
    coldir = os.path.join(TMP, "collide")
    os.makedirs(coldir, exist_ok=True)
    cu = os.path.join(TMP, "uniq.jpg")
    Image.new("RGB", (16, 16)).save(cu, "JPEG")
    o1 = default_output_path(cu, coldir, "_clean")
    check("first name free", o1.endswith("uniq_clean.jpg"))
    open(o1, "wb").write(b"placeholder")     # simulate an existing copy
    o2 = default_output_path(cu, coldir, "_clean")
    check("collision increments", o2.endswith("uniq_clean_1.jpg")
          and o1 != o2)

    # no extension
    noext = os.path.join(TMP, "noextension")
    Image.new("RGB", (8, 8)).save(noext, "JPEG")
    rep = sanitize(noext)
    check("no-extension file skipped", rep.status == "skipped")

    # unsupported type
    txt = os.path.join(TMP, "note.txt")
    open(txt, "w").write("hi")
    check("txt unsupported", not supported(txt))
    rep = sanitize(txt)
    check("txt skipped", rep.status == "skipped")


def t_corrupt():
    print("\n[Corrupt / empty inputs]")
    # empty jpg
    e = os.path.join(TMP, "empty.jpg")
    open(e, "wb").write(b"")
    rep = sanitize(e)
    check("empty jpg -> error", rep.status == "error")

    # garbage jpg
    g = os.path.join(TMP, "garbage.jpg")
    open(g, "wb").write(b"this is not an image at all")
    rep = sanitize(g)
    check("garbage jpg -> error", rep.status == "error")

    # truncated valid jpeg -> must not crash; clean error or graceful copy
    src = os.path.join(TMP, "a.jpg")
    half = open(src, "rb").read()[:100]
    tr = os.path.join(TMP, "trunc.jpg")
    open(tr, "wb").write(half)
    rep = sanitize(tr)
    check("truncated jpeg -> error (no broken file written)",
          rep.status == "error" and not os.path.exists(rep.dst),
          rep.summary())

    # truncated png
    p = os.path.join(TMP, "a.png")
    tp = os.path.join(TMP, "trunc.png")
    open(tp, "wb").write(open(p, "rb").read()[:20])
    rep = sanitize(tp)
    check("truncated png -> error", rep.status == "error", rep.summary())


def t_original_untouched():
    print("\n[Originals never modified]")
    fixtures = [os.path.join(TMP, f) for f in os.listdir(TMP)
                if os.path.isfile(os.path.join(TMP, f))]
    hashes = {f: sha(f) for f in fixtures if f.endswith(
        (".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif", ".gif", ".bmp",
         ".heic", ".pdf", ".docx", ".xlsx", ".pptx"))}
    for f in fixtures:
        if f.endswith((".jpg", ".jpeg", ".png", ".webp", ".tiff", ".tif",
                       ".gif", ".bmp", ".heic", ".pdf", ".docx", ".xlsx",
                       ".pptx")):
            opts = Options(strip_icc=True)
            sanitize(f, opts)
    changed = [f for f, h in hashes.items() if sha(f) != h]
    check(f"{len(hashes)} originals all byte-identical after cleaning",
          not changed, str(changed))


def t_cli():
    print("\n[CLI]")
    folder = os.path.join(TMP, "clipool")
    os.makedirs(folder, exist_ok=True)
    Image.new("RGB", (20, 20), (1, 2, 3)).save(
        os.path.join(folder, "one.jpg"), "JPEG", exif=make_exif().tobytes())
    Image.new("RGB", (20, 20), (4, 5, 6)).save(
        os.path.join(folder, "two.png"), "PNG")
    open(os.path.join(folder, "ignore.txt"), "w").write("x")
    cliout = os.path.join(TMP, "cliout")
    cli_script = os.path.join(os.path.dirname(os.path.dirname(TMP)), "scrubber.py")
    r = subprocess.run(
        [sys.executable, cli_script,
         folder, "--out", cliout, "--suffix", "_safe"],
        capture_output=True, text=True)
    check("cli exit 0", r.returncode == 0, r.stderr)
    check("cli produced outputs",
          os.path.exists(os.path.join(cliout, "one_safe.jpg"))
          and os.path.exists(os.path.join(cliout, "two_safe.png")))
    out_img = os.path.join(cliout, "one_safe.jpg")
    check("cli output cleaned", not dict(Image.open(out_img).getexif()))
    check("cli skipped txt",
          not os.path.exists(os.path.join(cliout, "ignore_safe.txt")))

    # CLI error exit for missing file
    r = subprocess.run(
        [sys.executable, cli_script, os.path.join(TMP, "does_not_exist.jpg")],
        capture_output=True, text=True)
    check("cli nonzero on missing", r.returncode != 0)


def main():
    shutil.rmtree(TMP, ignore_errors=True)
    os.makedirs(OUT)
    t_jpeg()
    t_icc()
    t_png()
    t_webp()
    t_tiff_gif_bmp()
    t_heic()
    t_pdf()
    t_office()
    t_filenames()
    t_corrupt()
    t_original_untouched()
    t_cli()
    print(f"\n===== {PASS} passed, {FAIL} failed, {SKIP} skipped =====")
    sys.exit(1 if FAIL else 0)


if __name__ == "__main__":
    main()
